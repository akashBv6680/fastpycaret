# utils/presentation.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
import os
import tempfile
import pandas as pd

def create_presentation(df: pd.DataFrame, target_column: str, task_type: str) -> str:
    """Creates a simple PowerPoint presentation summarizing the dataset and task."""
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "📊 AutoML Data Science Summary"
    slide.placeholders[1].text = f"Task: {task_type.title()} | Target: {target_column}"

    # Slide 2: Dataset Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "📁 Dataset Overview"
    content = f"- Rows: {df.shape[0]}\n- Columns: {df.shape[1]}\n- Target Column: {target_column}"
    textbox = slide.shapes.placeholders[1]
    textbox.text = content

    # Slide 3: Sample Data
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "🔍 Sample Data"
    textbox = slide.shapes.placeholders[1]
    sample_text = df.head(5).to_string(index=False)
    textbox.text = sample_text

    # Save presentation
    temp_dir = tempfile.mkdtemp()
    pptx_path = os.path.join(temp_dir, "summary_presentation.pptx")
    prs.save(pptx_path)
    return pptx_path
