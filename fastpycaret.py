# === app.py ===
import streamlit as st
import pandas as pd
import os
from utils.eda_generator import generate_eda_report
from utils.model_trainer import train_model_and_report
from utils.report_sender import send_email_with_attachments, is_valid_email
from utils.presentation import create_presentation
from utils.email_ai_responder import fetch_latest_email, generate_ai_reply

st.set_page_config(page_title="🤖 Agentic Data Science AI", layout="wide")
st.title("🤖 Agentic Data Science Automation")

# --- Dataset Upload ---
st.header("📁 Upload Dataset")
file = st.file_uploader("Upload a CSV, Excel, or JSON file", type=["csv", "xlsx", "xls", "json"])

def load_data(uploaded_file):
    if uploaded_file.name.endswith(".csv"):
        return pd.read_csv(uploaded_file)
    elif uploaded_file.name.endswith((".xls", ".xlsx")):
        return pd.read_excel(uploaded_file)
    elif uploaded_file.name.endswith(".json"):
        return pd.read_json(uploaded_file)

# --- File Processing ---
if file:
    df = load_data(file)
    st.session_state.df = df
    st.success("✅ Dataset loaded successfully.")
    st.dataframe(df.head())

    target_column = st.selectbox("🎯 Select Target Column for Modeling", df.columns)

    # Run EDA + Modeling
    if st.button("🚀 Run Full Data Science Pipeline"):
        with st.spinner("Generating EDA Report..."):
            eda_pdf = generate_eda_report(df)
            st.session_state.eda_pdf = eda_pdf
            st.success("EDA Report Generated ✅")

        with st.spinner("Training ML Model and Generating Report..."):
            model_pdf, task_type = train_model_and_report(df, target_column)
            st.session_state.model_pdf = model_pdf
            st.session_state.task_type = task_type
            st.success("Model Report Generated ✅")

        with st.spinner("Generating PowerPoint Presentation..."):
            pptx_path = create_presentation(df, target_column, task_type)
            st.session_state.presentation = pptx_path
            st.success("Presentation Created ✅")

# --- Email Reports ---
st.markdown("---")
st.header("📧 Share Reports with Client")
recipient_email = st.text_input("Enter Client's Email Address:")

if st.button("📤 Send Reports to Client"):
    if not recipient_email:
        st.warning("Please enter a client email address.")
    elif not is_valid_email(recipient_email):
        st.warning("⚠️ Invalid email address format. Please check and try again.")
    else:
        attachments = []
        for key in ['eda_pdf', 'model_pdf', 'presentation']:
            if key in st.session_state:
                attachments.append(st.session_state[key])

        subject = "Your Data Science Reports Are Ready!"
        body = f"""
Dear Client,

Your data science analysis is complete. Please find attached:

1. EDA Report
2. Model Performance Report
3. PowerPoint Summary

Task Type: {st.session_state.get('task_type', 'N/A').title()}

Best regards,
Your AutoML Agent
"""
        if send_email_with_attachments(recipient_email, subject, body, attachments):
            st.success("✅ Reports sent successfully!")
        else:
            st.error("❌ Failed to send email. Check server logs.")

# --- AI Email Responder ---
st.markdown("---")
st.header("📬 Automated Email Responder")
if st.button("📥 Check & Auto-Reply to Latest Unread Email"):
    from_email, subject, body = fetch_latest_email()
    if from_email:
        st.markdown(f"**From:** `{from_email}`")
        st.markdown(f"**Subject:** `{subject}`")
        st.text_area("📨 Message", body, height=200)

        with st.spinner("Generating reply..."):
            reply = generate_ai_reply(body)
        st.text_area("🤖 AI Reply", reply, height=200)

        if st.button("Send AI Reply"):
            if send_email_with_attachments(from_email, f"RE: {subject}", reply):
                st.success("✅ Reply sent!")
            else:
                st.error("❌ Failed to send reply.")
    else:
        st.info("📭 No new unread emails found.")

st.markdown("---")
st.markdown("Made with ❤️ by your Agentic AI. Powered by PyCaret + Streamlit")

# === utils/eda_generator.py ===
def generate_eda_report(df, output_path="eda_report.pdf"):
    import matplotlib.pyplot as plt
    from fpdf import FPDF
    import seaborn as sns
    import tempfile

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=14)
    pdf.cell(200, 10, txt="Exploratory Data Analysis Report", ln=1, align='C')

    stats = df.describe().round(2).to_string()
    pdf.set_font("Courier", size=10)
    for line in stats.splitlines():
        pdf.cell(0, 5, txt=line, ln=1)

    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmpfile:
        plt.figure(figsize=(10, 5))
        sns.heatmap(df.corr(numeric_only=True), annot=True, fmt=".2f", cmap="coolwarm")
        plt.title("Correlation Heatmap")
        plt.tight_layout()
        plt.savefig(tmpfile.name)
        pdf.image(tmpfile.name, w=180)

    pdf.output(output_path)
    return output_path

# === utils/model_trainer.py ===
def train_model_and_report(df, target, output_path="model_report.pdf"):
    from pycaret.classification import setup, compare_models, pull
    from fpdf import FPDF

    s = setup(df, target=target, silent=True, session_id=123)
    best_model = compare_models()
    report = pull()

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="Model Report", ln=1, align='C')
    for col in report.columns:
        pdf.cell(0, 10, txt=f"{col}: {report[col][0]}", ln=1)
    pdf.output(output_path)

    return output_path, s._get_config('ml_usecase')

# === utils/report_sender.py ===
def send_email_with_attachments(to_email, subject, body, attachments=[], from_email=os.getenv("EMAIL_ADDRESS"), password=os.getenv("EMAIL_PASSWORD")):
    import smtplib
    from email.mime.multipart import MIMEMultipart
    from email.mime.text import MIMEText
    from email.mime.base import MIMEBase
    from email import encoders

    msg = MIMEMultipart()
    msg['From'] = from_email
    msg['To'] = to_email
    msg['Subject'] = subject
    msg.attach(MIMEText(body, 'plain'))

    for file_path in attachments:
        with open(file_path, 'rb') as f:
            part = MIMEBase('application', 'octet-stream')
            part.set_payload(f.read())
            encoders.encode_base64(part)
            part.add_header('Content-Disposition', f'attachment; filename="{os.path.basename(file_path)}"')
            msg.attach(part)

    try:
        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
            server.login(from_email, password)
            server.send_message(msg)
        return True
    except Exception as e:
        print(e)
        return False

def is_valid_email(email):
    from email_validator import validate_email, EmailNotValidError
    try:
        validate_email(email)
        return True
    except EmailNotValidError:
        return False

# === utils/presentation.py ===
def create_presentation(df, target, task_type, output="summary.pptx"):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Data Science Summary"
    slide.placeholders[1].text = f"Target: {target}\nTask: {task_type}"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Dataset Overview"
    slide2.placeholders[1].text = str(df.describe().round(2))[:1000]

    prs.save(output)
    return output

# === utils/email_ai_responder.py ===
def fetch_latest_email():
    import imaplib
    import email
    from email.header import decode_header
    import os

    user = os.getenv("EMAIL_ADDRESS")
    password = os.getenv("EMAIL_PASSWORD")

    imap = imaplib.IMAP4_SSL("imap.gmail.com")
    imap.login(user, password)
    imap.select("inbox")
    status, messages = imap.search(None, "UNSEEN")

    messages = messages[0].split()
    if not messages:
        return None, None, None

    latest = messages[-1]
    _, msg_data = imap.fetch(latest, "(RFC822)")
    msg = email.message_from_bytes(msg_data[0][1])
    subject = decode_header(msg["Subject"])[0][0]
    if isinstance(subject, bytes):
        subject = subject.decode()
    from_email = msg.get("From")
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body = part.get_payload(decode=True).decode()
                break
    else:
        body = msg.get_payload(decode=True).decode()

    imap.logout()
    return from_email, subject, body

def generate_ai_reply(prompt):
    import together
    together.api_key = os.getenv("TOGETHER_API_KEY")
    response = together.Complete.create(
        prompt=f"Reply politely and informatively to this email:\n\n{prompt}",
        model="togethercomputer/llama-2-70b-chat",
        max_tokens=300
    )
    return response['output']['choices'][0]['text'].strip()
